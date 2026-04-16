"""简历解析服务：支持 PDF/DOC/DOCX/PPT/PPTX 格式"""

import json
import logging
import subprocess
import tempfile
from pathlib import Path

import fitz  # PyMuPDF
from sqlalchemy.orm import Session

from app.config import settings
from app.models.resume_models import ResumeFile, Resume
from app.services.ai import extract_structured_info

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str) -> str:
    """使用 PyMuPDF 从 PDF 提取纯文本"""
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text


def extract_text_from_docx(file_path: str) -> str:
    """使用 python-docx 从 DOCX 提取纯文本（含段落和表格）

    表格文本使用 [Table N|Row M] 标记保留结构信息，帮助 AI 识别工作经历、
    联系方式等表格区块。
    """
    from docx import Document
    doc = Document(file_path)
    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)
    # 提取表格中的文本，保留行列结构标记
    for t_idx, table in enumerate(doc.tables):
        for r_idx, row in enumerate(table.rows):
            cells = []
            for c_idx, cell in enumerate(row.cells):
                if cell.text.strip():
                    cells.append(cell.text.strip())
            if cells:
                row_text = " | ".join(cells)
                # 单列表格（智联招聘模板常见）：标记 Table/Row 帮助 AI 识别结构
                if len(table.columns) == 1:
                    texts.append(f"[表格{t_idx + 1} 行{r_idx + 1}] {row_text}")
                else:
                    texts.append(f"[表格{t_idx + 1}] {row_text}")
    return "\n".join(texts)


def extract_text_from_doc(file_path: str) -> str:
    """使用 libreoffice 将 DOC 转为 DOCX，再提取文本"""
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", tmpdir, file_path],
                capture_output=True, timeout=60,
            )
            if result.returncode != 0:
                logger.warning("libreoffice 转换失败: %s, stderr: %s", file_path, result.stderr.decode())
                return ""
            docx_path = str(Path(tmpdir) / Path(file_path).with_suffix(".docx").name)
            if not Path(docx_path).exists():
                logger.warning("libreoffice 转换后文件不存在: %s", docx_path)
                return ""
            return extract_text_from_docx(docx_path)
    except FileNotFoundError:
        logger.warning("libreoffice 未安装，无法解析 .doc 文件: %s", file_path)
        return ""
    except subprocess.TimeoutExpired:
        logger.warning("libreoffice 转换超时: %s", file_path)
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """使用 python-pptx 从 PPTX 提取纯文本"""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
    return "\n".join(texts)


def extract_text_from_ppt(file_path: str) -> str:
    """使用 libreoffice 将 PPT 转为 PPTX，再提取文本"""
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pptx", "--outdir", tmpdir, file_path],
                capture_output=True, timeout=60,
            )
            if result.returncode != 0:
                logger.warning("libreoffice 转换失败: %s, stderr: %s", file_path, result.stderr.decode())
                return ""
            pptx_path = str(Path(tmpdir) / Path(file_path).with_suffix(".pptx").name)
            if not Path(pptx_path).exists():
                logger.warning("libreoffice 转换后文件不存在: %s", pptx_path)
                return ""
            return extract_text_from_pptx(pptx_path)
    except FileNotFoundError:
        logger.warning("libreoffice 未安装，无法解析 .ppt 文件: %s", file_path)
        return ""
    except subprocess.TimeoutExpired:
        logger.warning("libreoffice 转换超时: %s", file_path)
        return ""


def extract_text_by_extension(file_path: str) -> str:
    """根据文件扩展名路由到对应的文本提取函数"""
    ext = Path(file_path).suffix.lower()
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".doc": extract_text_from_doc,
        ".pptx": extract_text_from_pptx,
        ".ppt": extract_text_from_ppt,
    }
    extractor = extractors.get(ext)
    if not extractor:
        logger.warning("不支持的文件格式: %s (%s)", ext, file_path)
        return ""
    try:
        return extractor(file_path)
    except Exception:
        logger.exception("文本提取失败: %s (%s)", ext, file_path)
        return ""


def parse_resume(db: Session, resume_file: ResumeFile) -> bool:
    """
    解析单份简历，将结构化数据写入数据库。

    支持 PDF/DOC/DOCX/PPT/PPTX 格式。
    """
    try:
        # 1. 根据扩展名提取文本
        raw_text = extract_text_by_extension(resume_file.file_path)
        if not raw_text.strip():
            logger.warning("简历内容为空: %s", resume_file.file_path)
            return False

        # 2. 调用 AI 提取结构化信息
        structured = extract_structured_info(raw_text)
        if structured is None:
            logger.error("AI 信息提取失败: %s", resume_file.file_path)
            return False

        # 3. 幂等写入：先删除该 file_id 的旧记录，再插入新的
        db.query(Resume).filter(Resume.file_id == resume_file.id).delete()

        resume = Resume(
            file_id=resume_file.id,
            name=structured.get("name"),
            email=structured.get("email"),
            phone=structured.get("phone"),
            current_title=structured.get("current_title"),
            years_exp=structured.get("years_exp"),
            education_json=json.dumps(structured.get("education", []), ensure_ascii=False),
            skills_json=json.dumps(structured.get("skills", []), ensure_ascii=False),
            work_exp_json=json.dumps(structured.get("work_experience", []), ensure_ascii=False),
            summary_text=structured.get("summary_text"),
        )
        db.add(resume)

        # 4. 更新文件状态
        resume_file.status = "done"
        resume_file.processed_at = __import__("datetime").datetime.now()

        db.commit()

        # 5. 建立向量索引
        try:
            from app.services.vector import add_resume_embedding, build_search_text
            search_text = build_search_text(resume)
            add_resume_embedding(resume.id, search_text)
        except Exception:
            logger.warning("向量索引建立失败（不影响主流程）", exc_info=True)

        db.commit()
        logger.info("简历解析成功: %s (%s)", structured.get("name", "Unknown"), resume_file.file_path)
        return True

    except Exception as e:
        db.rollback()
        resume_file.status = "failed"
        db.commit()
        logger.exception("简历解析失败: %s - %s", resume_file.file_path, e)
        return False
